"""
Document loaders: Load files from disk into Document objects.

Supports a plugin system for custom file type handlers::

    from ravenrag.loaders import register_loader

    @register_loader(".pdf")
    def load_pdf(path, metadata=None):
        ...  # return Document
"""

import logging
from pathlib import Path
from typing import Callable, Dict, List, Optional

from .index import Document

logger = logging.getLogger(__name__)

# Plugin registry: extension -> loader function
_loader_registry: Dict[str, Callable] = {}


def register_loader(extension: str):
    """Register a custom loader for a file extension.

    The decorated function receives (path: str, metadata: dict | None)
    and must return a Document.

    Example::

        @register_loader(".pdf")
        def load_pdf(path, metadata=None):
            text = extract_text_from_pdf(path)
            return Document(text=text, metadata={"source": path, **(metadata or {})})
    """

    def decorator(func: Callable) -> Callable:
        ext = extension if extension.startswith(".") else f".{extension}"
        _loader_registry[ext] = func
        logger.debug("Registered loader for %s: %s", ext, func.__name__)
        return func

    return decorator


def get_registered_extensions() -> List[str]:
    """Return list of file extensions with registered loaders."""
    return list(_loader_registry.keys())


def load_text(path: str, metadata: Optional[Dict] = None) -> Document:
    """Load a single text file as a Document."""
    file_path = Path(path).resolve()
    if not file_path.is_file():
        raise FileNotFoundError(f"File not found: {path}")

    text = file_path.read_text(encoding="utf-8")
    doc_metadata = {"source": str(file_path), "filename": file_path.name}
    if metadata:
        doc_metadata.update(metadata)

    return Document(text=text, metadata=doc_metadata)


def load_directory(
    path: str,
    glob: str = "**/*.txt",
    metadata: Optional[Dict] = None,
) -> List[Document]:
    """Load all matching files from a directory as Documents.

    Supports: .txt, .md, .py, .json, .csv, .html, .xml, .yaml, .yml, .toml
    Custom file types can be handled via ``register_loader()``.
    """
    dir_path = Path(path).resolve()
    if not dir_path.is_dir():
        raise NotADirectoryError(f"Not a directory: {path}")

    documents: List[Document] = []
    for file_path in sorted(dir_path.glob(glob)):
        if not file_path.is_file():
            continue

        # Path traversal protection: resolve symlinks and ensure within target
        resolved = file_path.resolve()
        try:
            resolved.relative_to(dir_path)
        except ValueError:
            logger.warning("Skipping %s: outside target directory (symlink?)", file_path)
            continue

        # Check for registered plugin loader
        if file_path.suffix in _loader_registry:
            try:
                doc = _loader_registry[file_path.suffix](str(file_path), metadata)
                documents.append(doc)
                continue
            except Exception as e:
                logger.warning("Plugin loader failed for %s: %s", file_path, e)
                continue

        # Default: read as text
        try:
            text = file_path.read_text(encoding="utf-8")
        except (UnicodeDecodeError, PermissionError) as e:
            logger.warning("Skipping %s: %s", file_path, e)
            continue

        doc_metadata = {
            "source": str(file_path),
            "filename": file_path.name,
            "extension": file_path.suffix,
        }
        if metadata:
            doc_metadata.update(metadata)

        documents.append(Document(text=text, metadata=doc_metadata))

    return documents


# ---------------------------------------------------------------------------
# Built-in file loaders (optional dependencies)
# ---------------------------------------------------------------------------


def _try_register_builtin_loaders() -> None:
    """Auto-register loaders for PDF, DOCX, HTML, and Markdown if deps are available."""

    # PDF via pymupdf4llm
    try:
        import pymupdf4llm  # noqa: F401

        @register_loader(".pdf")
        def _load_pdf(path: str, metadata: Optional[Dict] = None) -> Document:
            text = pymupdf4llm.to_markdown(path)
            file_path = Path(path).resolve()
            doc_metadata = {"source": str(file_path), "filename": file_path.name, "extension": ".pdf"}
            if metadata:
                doc_metadata.update(metadata)
            return Document(text=text, metadata=doc_metadata)

    except ImportError:
        pass

    # DOCX via python-docx
    try:
        import docx as _docx  # noqa: F401

        @register_loader(".docx")
        def _load_docx(path: str, metadata: Optional[Dict] = None) -> Document:
            doc = _docx.Document(path)
            paragraphs = [p.text for p in doc.paragraphs if p.text.strip()]
            text = "\n\n".join(paragraphs)
            file_path = Path(path).resolve()
            doc_metadata = {"source": str(file_path), "filename": file_path.name, "extension": ".docx"}
            if metadata:
                doc_metadata.update(metadata)
            return Document(text=text, metadata=doc_metadata)

    except ImportError:
        pass

    # HTML via beautifulsoup4
    try:
        from bs4 import BeautifulSoup  # noqa: F401

        @register_loader(".html")
        @register_loader(".htm")
        def _load_html(path: str, metadata: Optional[Dict] = None) -> Document:
            file_path = Path(path).resolve()
            raw = file_path.read_text(encoding="utf-8")
            soup = BeautifulSoup(raw, "html.parser")
            # Remove script and style elements
            for tag in soup(["script", "style", "nav", "footer", "header"]):
                tag.decompose()
            text = soup.get_text(separator="\n", strip=True)
            doc_metadata = {"source": str(file_path), "filename": file_path.name, "extension": file_path.suffix}
            if metadata:
                doc_metadata.update(metadata)
            return Document(text=text, metadata=doc_metadata)

    except ImportError:
        pass

    # CSV via stdlib (no extra deps)
    @register_loader(".csv")
    def _load_csv(path: str, metadata: Optional[Dict] = None) -> Document:
        import csv

        file_path = Path(path).resolve()
        with file_path.open(encoding="utf-8") as f:
            reader = csv.reader(f)
            rows = list(reader)
        text = "\n".join(", ".join(row) for row in rows)
        doc_metadata: Dict = {"source": str(file_path), "filename": file_path.name, "extension": ".csv"}
        if metadata:
            doc_metadata.update(metadata)
        return Document(text=text, metadata=doc_metadata)

    # RTF via striprtf
    try:
        from striprtf.striprtf import rtf_to_text  # noqa: F401

        @register_loader(".rtf")
        def _load_rtf(path: str, metadata: Optional[Dict] = None) -> Document:
            file_path = Path(path).resolve()
            raw = file_path.read_text(encoding="utf-8")
            text = rtf_to_text(raw)
            doc_metadata: Dict = {"source": str(file_path), "filename": file_path.name, "extension": ".rtf"}
            if metadata:
                doc_metadata.update(metadata)
            return Document(text=text, metadata=doc_metadata)

    except ImportError:
        pass

    # PPTX via python-pptx
    try:
        from pptx import Presentation  # noqa: F401

        @register_loader(".pptx")
        def _load_pptx(path: str, metadata: Optional[Dict] = None) -> Document:
            prs = Presentation(path)
            texts = []
            for slide in prs.slides:
                for shape in slide.shapes:
                    if shape.has_text_frame:
                        for paragraph in shape.text_frame.paragraphs:
                            text = paragraph.text.strip()
                            if text:
                                texts.append(text)
            file_path = Path(path).resolve()
            doc_metadata: Dict = {"source": str(file_path), "filename": file_path.name, "extension": ".pptx"}
            if metadata:
                doc_metadata.update(metadata)
            return Document(text="\n\n".join(texts), metadata=doc_metadata)

    except ImportError:
        pass

    # XLSX via openpyxl
    try:
        from openpyxl import load_workbook  # noqa: F401

        @register_loader(".xlsx")
        def _load_xlsx(path: str, metadata: Optional[Dict] = None) -> Document:
            wb = load_workbook(path, read_only=True, data_only=True)
            texts = []
            for sheet in wb.sheetnames:
                ws = wb[sheet]
                for row in ws.iter_rows(values_only=True):
                    cells = [str(c) if c is not None else "" for c in row]
                    line = ", ".join(cells).strip(", ")
                    if line:
                        texts.append(line)
            wb.close()
            file_path = Path(path).resolve()
            doc_metadata: Dict = {
                "source": str(file_path),
                "filename": file_path.name,
                "extension": ".xlsx",
            }
            if metadata:
                doc_metadata.update(metadata)
            return Document(text="\n".join(texts), metadata=doc_metadata)

    except ImportError:
        pass

    # Markdown with frontmatter parsing (no extra deps)
    @register_loader(".md")
    @register_loader(".markdown")
    def _load_markdown(path: str, metadata: Optional[Dict] = None) -> Document:
        file_path = Path(path).resolve()
        raw = file_path.read_text(encoding="utf-8")
        doc_metadata: Dict = {"source": str(file_path), "filename": file_path.name, "extension": ".md"}

        # Parse YAML frontmatter
        if raw.startswith("---"):
            parts = raw.split("---", 2)
            if len(parts) >= 3:
                frontmatter_text = parts[1].strip()
                raw = parts[2].strip()
                # Simple YAML key: value parsing (no dependency)
                for line in frontmatter_text.splitlines():
                    if ":" in line:
                        key, _, value = line.partition(":")
                        key = key.strip()
                        value = value.strip().strip("\"'")
                        if key and value:
                            doc_metadata[key] = value

        if metadata:
            doc_metadata.update(metadata)
        return Document(text=raw, metadata=doc_metadata)


# Auto-register on import
_try_register_builtin_loaders()
